import zipfile
from io import BytesIO
from pathlib import Path

# Per-entry uncompressed size cap when unpacking a zip, so a small/malicious
# archive can't decompress into a huge amount of text in memory (zip bomb).
_ZIP_ENTRY_MAX_MB = 20


def _load_pdf(raw_bytes: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(raw_bytes))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _load_docx(raw_bytes: bytes) -> str:
    import docx

    document = docx.Document(BytesIO(raw_bytes))
    return "\n\n".join(p.text for p in document.paragraphs)


def _load_text(raw_bytes: bytes) -> str:
    return raw_bytes.decode("utf-8", errors="replace")


def _load_xlsx(raw_bytes: bytes) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(BytesIO(raw_bytes), data_only=True, read_only=True)
    sheets = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _load_pptx(raw_bytes: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(BytesIO(raw_bytes))
    slides = []
    for i, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text for cell in row.cells))
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _load_zip(raw_bytes: bytes) -> str:
    """Unpack a zip and extract text from every supported file inside it
    (nested zips are skipped, not recursed into, to avoid zip-bomb-style
    blowups), concatenating each file's content under a filename header so
    the model can tell which part of the answer came from which file."""
    try:
        archive = zipfile.ZipFile(BytesIO(raw_bytes))
    except zipfile.BadZipFile:
        raise ValueError("Not a valid zip archive")

    sections = []
    with archive:
        for info in archive.infolist():
            if info.is_dir():
                continue

            name = info.filename
            base_name = Path(name).name
            if base_name.startswith(".") or "__MACOSX" in name:
                continue

            ext = Path(name).suffix.lower()
            loader = _LOADERS.get(ext)
            if loader is None or ext == ".zip":
                continue
            if info.file_size > _ZIP_ENTRY_MAX_MB * 1024 * 1024:
                continue

            try:
                text = loader(archive.read(info))
            except Exception:
                continue

            if text and text.strip():
                sections.append(f"=== {name} ===\n\n{text.strip()}")

    if not sections:
        raise ValueError(
            "No readable PDF, DOCX, TXT, MD, CSV, JSON, XLSX or PPTX files found inside the zip"
        )

    return "\n\n".join(sections)


_LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".txt": _load_text,
    ".md": _load_text,
    ".csv": _load_text,
    ".json": _load_text,
    ".xlsx": _load_xlsx,
    ".pptx": _load_pptx,
    ".zip": _load_zip,
}


def load_document(filename: str, raw_bytes: bytes) -> str:
    ext = Path(filename).suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported file type: {ext or '(none)'}")
    return loader(raw_bytes)
